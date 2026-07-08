"""
PowerPoint (pptx) -> PDF feature.
Route: POST /api/pptx-to-pdf
"""
import io
from flask import Blueprint, request, send_file, jsonify
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

pptx_to_pdf_bp = Blueprint('pptx_to_pdf', __name__)


@pptx_to_pdf_bp.route('/api/pptx-to-pdf', methods=['POST'])
def pptx_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No presentation uploaded"}), 400
    file = request.files['file']
    try:
        prs = Presentation(file.stream)
        pdf_buffer = io.BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=letter)

        for i, slide in enumerate(prs.slides):
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, 750, f"Slide Layer Structure Asset {i+1}")
            y_pos = 700
            c.setFont("Helvetica", 10)

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    c.drawString(50, y_pos, shape.text[:90])
                    y_pos -= 20
            c.showPage()

        c.save()
        pdf_buffer.seek(0)
        return send_file(pdf_buffer, as_attachment=True, download_name="presentation_converted.pdf", mimetype="application/pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
