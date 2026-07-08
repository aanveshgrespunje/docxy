"""
PDF -> PowerPoint (pptx) feature.
Route: POST /api/pdf-to-pptx
"""
import io
from flask import Blueprint, request, send_file, jsonify
import pdfplumber
from pptx import Presentation
from pptx.util import Inches

pdf_to_pptx_bp = Blueprint('pdf_to_pptx', __name__)


@pdf_to_pptx_bp.route('/api/pdf-to-pptx', methods=['POST'])
def pdf_to_pptx():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    file = request.files['file']
    try:
        pdf_bytes = file.read()
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                if text:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    tf.text = text

        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return send_file(pptx_buffer, as_attachment=True, download_name="reverse_parsed_slides.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
